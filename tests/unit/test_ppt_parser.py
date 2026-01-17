#!/usr/bin/env python3
"""
PPT解析器单元测试

测试场景:
- 空PPT（0页）
- 单页PPT
- 多页PPT（50+页）
- 复杂布局PPT
- 包含特殊字符的PPT
- 损坏的PPTX文件

作者: Claude Code
版本: v1.0
"""

import unittest
import sys
import os
import time
from pathlib import Path
from io import BytesIO

# 添加项目根目录到Python路径
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from src.modules.ppt_parser import extract_text_from_ppt, extract_tables_from_ppt


class TestPPTParser(unittest.TestCase):
    """PPT解析器测试类"""

    def setUp(self):
        """测试前准备"""
        self.test_ppt_dir = Path(project_root / "tests" / "test_ppts")
        self.temp_dir = Path(project_root / "tmp" / "test_output")
        self.temp_dir.mkdir(parents=True, exist_ok=True)

    def tearDown(self):
        """测试后清理"""
        pass

    def test_extract_text_from_ppt(self):
        """测试文本提取功能"""
        ppt_path = self.test_ppt_dir / "sample_ai_course.pptx"
        result = extract_text_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertEqual(len(result), 3)

        for slide in result:
            self.assertIn('page', slide)
            self.assertIn('title', slide)
            self.assertIn('content', slide)
            self.assertIn('hyperlinks', slide)
            self.assertIn('shapes', slide)

    def test_extract_tables_from_ppt(self):
        """测试表格提取功能"""
        ppt_path = self.test_ppt_dir / "complex_sample.pptx"
        result = extract_tables_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertEqual(len(result), 1)

        table = result[0]
        self.assertEqual(table['slide_number'], 2)
        self.assertEqual(table['rows'], 3)
        self.assertEqual(table['cols'], 3)
        self.assertIsInstance(table['data'], list)

    def test_empty_ppt(self):
        """测试空PPT处理"""
        ppt_path = self.test_ppt_dir / "empty.pptx"
        result = extract_text_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertEqual(len(result), 0)

    def test_single_page_ppt(self):
        """测试单页PPT处理"""
        from pptx import Presentation

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "单页测试"
        slide.placeholders[1].text = "这是单页PPT的内容"

        ppt_path = self.temp_dir / "single_page.pptx"
        prs.save(str(ppt_path))

        result = extract_text_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertEqual(len(result), 1)
        self.assertEqual(result[0]['page'], 1)
        self.assertEqual(result[0]['title'], "单页测试")

    def test_multi_page_ppt(self):
        """测试多页PPT处理"""
        ppt_path = self.test_ppt_dir / "multi_pages.pptx"
        start_time = time.time()
        result = extract_text_from_ppt(str(ppt_path))
        end_time = time.time()

        # 验证结果
        self.assertIsInstance(result, list)
        self.assertEqual(len(result), 20)

        # 验证性能（20页应该在10秒内处理完成）
        processing_time = end_time - start_time
        self.assertLess(processing_time, 10.0,
                        f"处理时间过长: {processing_time:.2f}秒")

    def test_complex_layout_ppt(self):
        """测试复杂布局PPT"""
        ppt_path = self.test_ppt_dir / "complex_sample.pptx"
        result = extract_text_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        for slide in result:
            self.assertIn('shapes', slide)
            self.assertIsInstance(slide['shapes'], list)

    def test_special_characters(self):
        """测试特殊字符处理"""
        ppt_path = self.test_ppt_dir / "special_chars.pptx"
        result = extract_text_from_ppt(str(ppt_path))

        self.assertIsInstance(result, list)
        self.assertGreater(len(result), 0)

        # 验证特殊字符被正确处理
        all_text = ""
        for slide in result:
            all_text += slide['title'] + " "
            for content in slide['content']:
                all_text += content + " "

        self.assertIn("中文", all_text)
        self.assertIn("测试", all_text)

    def test_invalid_file(self):
        """测试无效文件处理"""
        invalid_path = "/nonexistent/path/file.pptx"

        with self.assertRaises(FileNotFoundError):
            extract_text_from_ppt(invalid_path)

    def test_corrupted_ppt(self):
        """测试损坏PPT处理"""
        corrupted_path = self.temp_dir / "corrupted.pptx"
        with open(corrupted_path, 'wb') as f:
            f.write(b"This is not a valid PPTX file")

        with self.assertRaises(Exception):
            extract_text_from_ppt(str(corrupted_path))

    def test_performance_large_ppt(self):
        """测试大文件性能"""
        from pptx import Presentation

        # 创建50页PPT
        prs = Presentation()
        for i in range(50):
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"第{i+1}页"
            content = slide.placeholders[1].text_frame
            content.text = f"这是第{i+1}页的大量文本内容。" * 10

        ppt_path = self.temp_dir / "large_ppt.pptx"
        prs.save(str(ppt_path))

        # 测试处理时间
        start_time = time.time()
        result = extract_text_from_ppt(str(ppt_path))
        end_time = time.time()

        processing_time = end_time - start_time

        # 验证结果
        self.assertEqual(len(result), 50)

        # 性能要求：50页应该在30秒内处理完成
        self.assertLess(processing_time, 30.0,
                        f"大文件处理时间过长: {processing_time:.2f}秒")

        print(f"\n大文件性能测试结果:")
        print(f"  页数: {len(result)}")
        print(f"  处理时间: {processing_time:.2f}秒")
        print(f"  平均每页: {processing_time/50*1000:.2f}毫秒")

    def test_memory_usage(self):
        """测试内存使用"""
        import psutil
        import os

        process = psutil.Process(os.getpid())
        mem_before = process.memory_info().rss / 1024 / 1024  # MB

        ppt_path = self.test_ppt_dir / "multi_pages.pptx"
        result = extract_text_from_ppt(str(ppt_path))

        mem_after = process.memory_info().rss / 1024 / 1024  # MB
        mem_increase = mem_after - mem_before

        # 内存增长应该小于500MB
        self.assertLess(mem_increase, 500,
                        f"内存使用过多: {mem_increase:.2f}MB")

        print(f"\n内存使用测试结果:")
        print(f"  处理前: {mem_before:.2f}MB")
        print(f"  处理后: {mem_after:.2f}MB")
        print(f"  增长: {mem_increase:.2f}MB")


if __name__ == "__main__":
    # 运行测试
    print("="*60)
    print("开始PPT解析器单元测试")
    print("="*60)

    unittest.main(verbosity=2)
