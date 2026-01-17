#!/usr/bin/env python3
"""
PPT解析模块 - 提取PPT文件中的文字、图片、布局信息

功能:
- 提取每页PPT的文字内容
- 提取图片和其他媒体元素
- 生成结构化数据
- 将PPT转换为图片序列

作者: Claude Code
版本: v1.0
"""

from pathlib import Path
from typing import List, Dict, Any, Optional


def extract_text_from_ppt(ppt_path: str) -> List[Dict[str, Any]]:
    """
    从PPT文件中提取文本内容

    Args:
        ppt_path: PPT文件路径

    Returns:
        List[Dict]: 包含每页文本内容的列表
            - page: 页码
            - title: 标题
            - content: 正文内容列表
            - hyperlinks: 超链接列表（URL和文本）
            - shapes: 元素位置信息列表

    Raises:
        FileNotFoundError: PPT文件不存在
        Exception: PPT解析失败
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 检查文件是否存在
    if not Path(ppt_path).exists():
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")

    slides_data = []

    try:
        # 打开PPT文件
        prs = Presentation(ppt_path)

        # 遍历每页幻灯片
        for page_num, slide in enumerate(prs.slides, 1):
            slide_text = {
                'page': page_num,
                'title': '',
                'content': [],
                'hyperlinks': [],
                'shapes': []
            }

            # 提取所有文本内容
            for shape_idx, shape in enumerate(slide.shapes):
                # 记录元素位置信息
                shape_info = {
                    'index': shape_idx,
                    'type': str(shape.shape_type),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }

                # 处理文本框
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()

                    # 提取超链接
                    hyperlinks = []
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run, 'hyperlink') and run.hyperlink:
                                    try:
                                        # 尝试不同的属性名
                                        url = ''
                                        if hasattr(run.hyperlink, 'address'):
                                            url = run.hyperlink.address
                                        elif hasattr(run.hyperlink, 'target'):
                                            url = run.hyperlink.target
                                        elif hasattr(run.hyperlink, 'rId'):
                                            # 如果是关系ID，需要解析
                                            url = f"<relationship:{run.hyperlink.rId}>"

                                        hyperlinks.append({
                                            'text': run.text,
                                            'url': url
                                        })
                                        shape_info['has_hyperlink'] = True
                                    except Exception as e:
                                        # 忽略超链接解析错误
                                        pass

                    if hyperlinks:
                        slide_text['hyperlinks'].extend(hyperlinks)

                    # 尝试识别标题
                    # 方案1：通过占位符类型判断
                    if shape.is_placeholder:
                        placeholder = shape.placeholder_format
                        if placeholder.type == 1:  # 标题占位符
                            if not slide_text['title']:
                                slide_text['title'] = text
                            shape_info['is_title'] = True
                            slide_text['shapes'].append(shape_info)
                            continue

                    # 方案2：通过位置和字体大小判断（标题通常在上方且字体较大）
                    if not slide_text['title'] and len(text) < 100:
                        # 如果是第一个文本且较短，可能是标题
                        slide_text['title'] = text
                        shape_info['is_title'] = True
                    else:
                        slide_text['content'].append(text)

                    slide_text['shapes'].append(shape_info)

                # 处理表格中的文本
                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_content = []
                    table = shape.table
                    for row in table.rows:
                        row_content = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_content.append(cell.text.strip())
                        if row_content:
                            table_content.append(" | ".join(row_content))
                    if table_content:
                        slide_text['content'].append(f"表格: {' ; '.join(table_content)}")
                    shape_info['is_table'] = True
                    slide_text['shapes'].append(shape_info)

                # 处理图片
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    shape_info['is_picture'] = True
                    slide_text['shapes'].append(shape_info)

                # 其他形状类型
                else:
                    slide_text['shapes'].append(shape_info)

            # 如果没有标题，从内容中提取第一个作为标题
            if not slide_text['title'] and slide_text['content']:
                slide_text['title'] = slide_text['content'][0]
                slide_text['content'] = slide_text['content'][1:]

            slides_data.append(slide_text)

        return slides_data

    except Exception as e:
        raise Exception(f"PPT解析失败: {str(e)}") from e


def convert_ppt_to_images(ppt_path: str, output_dir: str) -> List[str]:
    """
    将PPT幻灯片转换为图片

    Args:
        ppt_path: PPT文件路径
        output_dir: 输出目录

    Returns:
        List[str]: 生成的图片文件路径列表

    Raises:
        Exception: 转换失败
    """
    import subprocess
    import shutil
    from pathlib import Path

    # 检查文件是否存在
    if not Path(ppt_path).exists():
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")

    # 创建输出目录
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    try:
        # 方法：LibreOffice命令行转换
        # 检查LibreOffice是否可用
        libreoffice_path = None
        import shutil
        if shutil.which('libreoffice'):
            libreoffice_path = 'libreoffice'
        else:
            # 尝试常见安装位置
            common_paths = [
                '/Applications/LibreOffice.app/Contents/MacOS/soffice',
                '/usr/bin/libreoffice',
                '/usr/local/bin/libreoffice'
            ]
            for path in common_paths:
                if Path(path).exists():
                    libreoffice_path = path
                    break

        if not libreoffice_path:
            raise Exception("未找到LibreOffice，请安装或添加到PATH")

        # 步骤1: 使用LibreOffice将PPT转换为PDF
        pdf_path = Path(output_dir) / f"{Path(ppt_path).stem}.pdf"

        # 运行LibreOffice转换命令
        cmd_convert = [
            libreoffice_path,
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', str(output_dir),
            ppt_path
        ]

        print(f"正在转换PPT为PDF: {ppt_path}")
        result = subprocess.run(
            cmd_convert,
            capture_output=True,
            text=True,
            timeout=60
        )

        if result.returncode != 0:
            raise Exception(f"LibreOffice转换失败: {result.stderr}")

        # 步骤2: 使用pdftoppm将PDF转换为图片序列
        print(f"正在将PDF转换为图片: {pdf_path}")

        stem = Path(ppt_path).stem
        image_prefix = str(Path(output_dir) / f"{stem}_slide_")

        # 运行pdftoppm转换命令
        # 语法: pdftoppm [options] PDF-file [output-prefix]
        cmd_images = [
            'pdftoppm',
            '-png',  # 输出PNG格式
            '-f', '1',  # 第一页
            '-l', '1000',  # 最大页数
            str(pdf_path),  # 输入PDF文件
            image_prefix  # 输出文件前缀
        ]

        result = subprocess.run(
            cmd_images,
            capture_output=True,
            text=True,
            timeout=120
        )

        if result.returncode != 0:
            raise Exception(f"pdftoppm转换失败: {result.stderr}")

        # 步骤3: 收集生成的图片文件
        image_files = []

        # 查找所有生成的图片文件
        for img_file in sorted(Path(output_dir).glob(f"{stem}_slide_*.png")):
            image_files.append(str(img_file))

        # 清理临时PDF文件
        if pdf_path.exists():
            pdf_path.unlink()

        print(f"成功转换 {len(image_files)} 张图片")
        return image_files

    except subprocess.TimeoutExpired:
        raise Exception("转换超时，请检查文件大小和网络连接")
    except Exception as e:
        raise Exception(f"PPT转图片失败: {str(e)}") from e


def extract_images_from_ppt(ppt_path: str, output_dir: str) -> List[str]:
    """
    从PPT中提取所有图片

    Args:
        ppt_path: PPT文件路径
        output_dir: 图片输出目录

    Returns:
        List[str]: 提取的图片文件路径列表
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from PIL import Image
    import io
    import uuid
    from pathlib import Path

    # 检查文件是否存在
    if not Path(ppt_path).exists():
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")

    # 创建输出目录
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    extracted_images = []

    try:
        # 打开PPT文件
        prs = Presentation(ppt_path)

        # 遍历每页幻灯片
        for page_num, slide in enumerate(prs.slides, 1):
            for shape_idx, shape in enumerate(slide.shapes):
                # 检查是否为图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        # 获取图片数据
                        image = shape.image
                        image_bytes = image.blob

                        # 生成唯一文件名
                        image_ext = image.ext
                        if not image_ext.startswith('.'):
                            image_ext = '.' + image_ext

                        image_filename = f"slide_{page_num:03d}_img_{shape_idx:03d}_{uuid.uuid4().hex[:8]}{image_ext}"
                        image_path = Path(output_dir) / image_filename

                        # 保存图片
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)

                        extracted_images.append(str(image_path))
                        print(f"提取图片: {image_filename}")

                    except Exception as e:
                        print(f"提取第{page_num}页第{shape_idx}个图片失败: {e}")

        return extracted_images

    except Exception as e:
        raise Exception(f"PPT图片提取失败: {str(e)}") from e


def extract_tables_from_ppt(ppt_path: str) -> List[Dict[str, Any]]:
    """
    从PPT中提取表格数据

    Args:
        ppt_path: PPT文件路径

    Returns:
        List[Dict]: 表格数据列表
            - slide_number: 幻灯片编号
            - table_index: 表格索引
            - rows: 行数
            - cols: 列数
            - data: 二维数组表格数据
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # 检查文件是否存在
    if not Path(ppt_path).exists():
        raise FileNotFoundError(f"PPT文件不存在: {ppt_path}")

    tables_data = []

    try:
        # 打开PPT文件
        prs = Presentation(ppt_path)

        # 遍历每页幻灯片
        for page_num, slide in enumerate(prs.slides, 1):
            table_index = 0

            for shape in slide.shapes:
                # 检查是否为表格
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    rows = table.rows
                    cols = table.columns

                    # 提取表格数据
                    table_data = []
                    for row_idx, row in enumerate(rows):
                        row_data = []
                        for col_idx, cell in enumerate(row.cells):
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)
                        table_data.append(row_data)

                    # 记录表格信息
                    table_info = {
                        'slide_number': page_num,
                        'table_index': table_index,
                        'rows': len(rows),
                        'cols': len(cols),
                        'data': table_data
                    }

                    tables_data.append(table_info)
                    table_index += 1

                    print(f"提取表格: 第{page_num}页, 表格{table_index}, {len(rows)}行x{len(cols)}列")

        return tables_data

    except Exception as e:
        raise Exception(f"PPT表格提取失败: {str(e)}") from e


if __name__ == "__main__":
    # 测试代码
    print("PPT解析模块已加载")
    print("注意: 模块尚未实现具体功能")
